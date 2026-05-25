# generate_ppt.py — Genera version PPT de clase_espejo_tercera_respuesta.html
import pathlib, re
from html import unescape
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colores tema ─────────────────────────────────────────────────────────────
BG        = RGBColor(0x0f, 0x0f, 0x1a)
BG_LIGHT  = RGBColor(0xf4, 0xf1, 0xeb)
GOLD      = RGBColor(0xff, 0xd7, 0x00)
CYAN      = RGBColor(0x00, 0xb4, 0xd8)
WHITE     = RGBColor(0xff, 0xff, 0xff)
GRAY      = RGBColor(0x99, 0x99, 0xaa)
RED       = RGBColor(0xe0, 0x41, 0x3a)
T_COL     = RGBColor(0x4c, 0xd9, 0x64)
I_COL     = RGBColor(0xff, 0xa0, 0x20)
F_COL     = RGBColor(0xe0, 0x41, 0x3a)
DIM       = RGBColor(0x55, 0x55, 0x66)

W  = Inches(13.333)
H  = Inches(7.5)
M  = Inches(0.55)   # margin

# ── HTML helpers ─────────────────────────────────────────────────────────────
def strip_html(s):
    s = re.sub(r'<[^>]+>', ' ', s)
    s = unescape(s)
    s = re.sub(r'\s+', ' ', s).strip()
    return s

def first_match(pattern, text, group=1, default=''):
    m = re.search(pattern, text, re.DOTALL | re.IGNORECASE)
    return strip_html(m.group(group)) if m else default

def all_matches(pattern, text):
    return [strip_html(m) for m in re.findall(pattern, text, re.DOTALL | re.IGNORECASE)]

# ── Slide helpers ─────────────────────────────────────────────────────────────
def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, x, y, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tf = slide.shapes.add_textbox(x, y, w, h)
    tf.word_wrap = wrap
    p = tf.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = 'Calibri'
    return tf

def add_accent_bar(slide, color, x=M, y=Inches(0.08), w=Inches(12.2), h=Inches(0.03)):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, fill_color, alpha_pct=100, line_color=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

# ── Parse HTML ────────────────────────────────────────────────────────────────
html = pathlib.Path('clase_espejo_tercera_respuesta.html').read_text(encoding='utf-8')

# Split by slide comments
parts = re.split(r'(<!-- ============ SLIDE [^=]+=+ -->)', html)

slides_data = []
i = 0
while i < len(parts):
    chunk = parts[i]
    m = re.match(r'<!-- ============ SLIDE ([^:]+): ([^=]+) =', chunk)
    if m:
        sid   = m.group(1).strip()
        stitle = m.group(2).strip()
        body  = parts[i+1] if i+1 < len(parts) else ''
        slides_data.append((sid, stitle, body))
        i += 2
    else:
        i += 1

print(f'Slides encontrados en HTML: {len(slides_data)}')

# ── Build PPT ─────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

def make_slide(sid, stitle, body, is_light=False):
    slide = prs.slides.add_slide(blank)
    bg_color = BG_LIGHT if is_light else BG
    set_bg(slide, bg_color)

    text_color  = RGBColor(0x1a,0x1a,0x2e) if is_light else WHITE
    muted_color = RGBColor(0x55,0x55,0x66) if is_light else GRAY
    eyebrow_c   = RGBColor(0x00,0x88,0xaa) if is_light else CYAN

    # Top accent bar
    add_accent_bar(slide, GOLD if not is_light else RGBColor(0xc8,0x58,0x20))

    # Slide number badge (top right)
    add_text(slide, sid, W - Inches(1.0), Inches(0.12), Inches(0.9), Inches(0.35),
             size=14, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

    # Badge label (top left small)
    badge = first_match(r'class="badge">([^<]+)<', body)
    if badge:
        add_text(slide, badge.upper(), M, Inches(0.1), Inches(8), Inches(0.3),
                 size=7, color=DIM)

    # Part / subtitle
    part = first_match(r'class="part">([^<]+)<', body)
    if part:
        add_text(slide, part, M, Inches(0.38), Inches(11), Inches(0.35),
                 size=9, color=muted_color)

    # Eyebrow
    eyebrow = first_match(r'class="eyebrow"[^>]*>([^<]{10,})<', body)
    if not eyebrow:
        eyebrow = first_match(r'class="eyebrow">([^<]{10,})<', body)
    if eyebrow:
        add_text(slide, eyebrow, M, Inches(0.78), Inches(11.5), Inches(0.4),
                 size=8.5, color=eyebrow_c, italic=True)

    # Main h2 title
    h2_raw = re.search(r'<h2[^>]*>(.*?)</h2>', body, re.DOTALL)
    h2 = strip_html(h2_raw.group(1)) if h2_raw else stitle
    h2 = h2[:160]
    add_text(slide, h2, M, Inches(1.2), Inches(12.2), Inches(1.5),
             size=28, bold=True, color=text_color, wrap=True)

    # Key content: extract meaningful text blocks (> 40 chars, skip buttons/tiny snippets)
    body_texts = []

    # Try display-s headlines first
    for m2 in re.finditer(r'class="display-s[^"]*">([^<]{20,})<', body):
        t = strip_html(m2.group(1))
        if t and t not in body_texts: body_texts.append(t)

    # Metric hero numbers
    for m2 in re.finditer(r'class="stat-hero[^"]*">([^<]{2,})<', body):
        t = strip_html(m2.group(1))
        if t and t not in body_texts: body_texts.append('▸ ' + t)

    # Paragraphs with real content
    for m2 in re.finditer(r'<(?:p|div)[^>]*>\s*([^<]{50,})\s*</', body):
        t = strip_html(m2.group(1))
        if t and len(t) > 45 and t not in body_texts:
            body_texts.append(t)
        if len(body_texts) >= 6: break

    # Display body content
    y_start = Inches(2.75)
    max_items = 5
    for idx2, txt in enumerate(body_texts[:max_items]):
        txt_short = txt[:200] + ('…' if len(txt) > 200 else '')
        y = y_start + Inches(idx2 * 0.82)
        if y + Inches(0.75) > H - Inches(0.4): break
        # bullet dot
        add_text(slide, '▸', M, y, Inches(0.25), Inches(0.7),
                 size=9, color=GOLD)
        add_text(slide, txt_short, M + Inches(0.28), y, Inches(11.8), Inches(0.75),
                 size=10.5, color=text_color, wrap=True)

    # Bottom strip: slide title
    label = f'SLIDE {sid}  ·  {stitle}'
    add_text(slide, label, M, H - Inches(0.35), W - M*2, Inches(0.3),
             size=7, color=DIM)

# ── Portada ──────────────────────────────────────────────────────────────────
cover = prs.slides.add_slide(blank)
set_bg(cover, BG)
add_accent_bar(cover, GOLD)
add_rect(cover, M, Inches(1.5), Inches(12.2), Pt(2), GOLD)
add_text(cover, 'La Tercera Respuesta en Analítica de Datos',
         M, Inches(1.7), Inches(12.2), Inches(1.5),
         size=36, bold=True, color=WHITE, wrap=True)
add_text(cover, 'Cuando la Incertidumbre No Es Ignorancia: Lógica Neutrosófica y LLMs',
         M, Inches(3.2), Inches(12.2), Inches(0.8),
         size=16, color=CYAN, italic=True)
add_text(cover, 'Dr. Maikel Leyva Vázquez  ·  UNIR · UBE  ·  2026',
         M, Inches(4.1), Inches(12.2), Inches(0.5),
         size=12, color=GRAY)
add_text(cover, 'mleyvaz.github.io/clase_espejo_tercera_respuesta.html',
         M, Inches(4.7), Inches(12.2), Inches(0.4),
         size=9, color=DIM)

# ── Generar un slide PPT por cada slide HTML ──────────────────────────────────
light_slides = {'7e'}  # slides con fondo claro en HTML

for sid, stitle, body in slides_data:
    is_light = any(sid.strip() == ls for ls in light_slides) or \
               'class="slide light"' in body
    make_slide(sid.strip(), stitle.strip(), body, is_light)
    print('  PPT slide:', sid.strip().encode('ascii','replace').decode(), '-', stitle.strip()[:50].encode('ascii','replace').decode())

# ── Guardar ───────────────────────────────────────────────────────────────────
out = pathlib.Path(r'C:\Users\HP\Documents\ClaseEspejo_TerceraRespuesta.pptx')
prs.save(str(out))
print(f'\nPPT guardado: {out}')
print(f'Total slides PPT: {len(prs.slides)} (1 portada + {len(slides_data)} contenido)')
